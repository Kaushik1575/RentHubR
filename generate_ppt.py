import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)        # #0f172a (Slate 900)
    CARD_BG = RGBColor(30, 41, 59)        # #1e293b (Slate 800)
    CARD_BORDER = RGBColor(51, 65, 85)    # #334155
    ACCENT_BLUE = RGBColor(59, 130, 246)  # #3b82f6
    ACCENT_CYAN = RGBColor(56, 189, 248)  # #38bdf8
    ACCENT_GREEN = RGBColor(16, 185, 129) # #10b981
    ACCENT_PURPLE = RGBColor(168, 85, 247)# #a855f7
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(203, 213, 225)  # #cbd5e1

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

    def add_header(slide, title, category="RENTHUB • MINOR PROJECT"):
        # Category Badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN

        # Main Slide Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    # Hero card
    hero = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.0), Inches(10.9), Inches(5.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = CARD_BG
    hero.line.color.rgb = CARD_BORDER
    hero.line.width = Pt(1.5)

    # Title Text
    t_box = slide1.shapes.add_textbox(Inches(1.8), Inches(1.4), Inches(9.7), Inches(4.8))
    tf = t_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "ACADEMIC MINOR PROJECT PRESENTATION"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_CYAN
    p0.space_after = Pt(10)

    p1 = tf.add_paragraph()
    p1.text = "RentHub"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Automated Vehicle Rental & Fleet Lifecycle Management System"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_BLUE
    p2.space_after = Pt(28)

    p3 = tf.add_paragraph()
    p3.text = "Presented by: Student / Project Team"
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_WHITE

    p4 = tf.add_paragraph()
    p4.text = "Branch: Computer Science & Engineering"
    p4.font.size = Pt(12)
    p4.font.color.rgb = TEXT_MUTED

    p5 = tf.add_paragraph()
    p5.text = "Domain: Web Technologies, Cloud Databases & Enterprise Systems"
    p5.font.size = Pt(12)
    p5.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement & Objectives
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "Project Overview & Key Objectives")

    # Card 1: Problem Statement
    c1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = RGBColor(239, 68, 68)

    t1 = slide2.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.6))
    tf1 = t1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "Existing Challenges in Vehicle Rental"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(248, 113, 113)
    p.space_after = Pt(14)

    points1 = [
        "Manual record keeping and delayed booking confirmations.",
        "Lack of real-time vehicle availability across categories.",
        "Physical paper agreements causing compliance disputes.",
        "No centralized moderation for user violations or fleet status.",
        "Complex deposit reconciliation and pricing transparency issues."
    ]
    for pt in points1:
        p = tf1.add_paragraph()
        p.text = "• " + pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    # Card 2: Proposed Solution
    c2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = ACCENT_GREEN

    t2 = slide2.shapes.add_textbox(Inches(7.2), Inches(1.9), Inches(5.0), Inches(4.6))
    tf2 = t2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "The RentHub Solution"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(14)

    points2 = [
        "100% Cloud-connected User & Admin web applications.",
        "Live fleet discovery with hourly dynamic pricing calculators.",
        "Digital KYC verification and integrated legal Rental Agreement Deed.",
        "Centralized Admin Command Room for fleet & booking controls.",
        "Automated offer codes, support ticketing & audit logs."
    ]
    for pt in points2:
        p = tf2.add_paragraph()
        p.text = "✔ " + pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 3: System Architecture & Tech Stack
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "System Architecture & Technology Stack")

    cols = [
        ("Frontend Layer", "React.js + Vite", ACCENT_CYAN, [
            "Single Page Application (SPA)",
            "Vanilla CSS Design System",
            "Glassmorphic Dark UI Theme",
            "Mobile Responsive Layouts",
            "Interactive Modals & Charts"
        ]),
        ("Backend & API Layer", "Node.js + Express", ACCENT_BLUE, [
            "RESTful API Endpoints",
            "JWT Authentication Guard",
            "Bcrypt Password Hashing",
            "Role-Based Access (RBAC)",
            "Automated Price Engines"
        ]),
        ("Database & Cloud", "Supabase + PostgreSQL", ACCENT_PURPLE, [
            "PostgreSQL Relational DB",
            "Row-Level Security Policies",
            "Cloudinary Document CDN",
            "Real-time Schema Sync",
            "ACID Compliant Transactions"
        ])
    ]

    for i, (col_title, sub, col_color, items) in enumerate(cols):
        left = Inches(0.8 + i * 3.9)
        c = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(5.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col_color
        c.line.width = Pt(1.5)

        t = slide3.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.3), Inches(4.8))
        tf = t.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = col_title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = col_color

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_after = Pt(14)

        for item in items:
            pi = tf.add_paragraph()
            pi.text = "• " + item
            pi.font.size = Pt(12)
            pi.font.color.rgb = TEXT_MUTED
            pi.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 4: User Module - Discovery, Catalog & Schedule
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "User Module: Vehicle Catalog & Discovery Engine", "USER EXPERIENCE LAYER")

    user_cards = [
        ("Fleet Categorization", "Bikes, Cars & Scooters", ACCENT_CYAN, [
            "Filter fleet by vehicle type (Bikes, Cars, Scooters).",
            "Inspect engine specs, fuel type, and hourly tariff.",
            "Live badge indicating available vs occupied status."
        ]),
        ("Trip Schedule & Tariff", "Dynamic Pricing Calculator", ACCENT_BLUE, [
            "Select pickup date, drop date, and duration.",
            "Real-time total fare estimation based on hourly rate.",
            "Transparent pricing with no hidden charges."
        ]),
        ("Offers & Promo Codes", "Instant Discount Reductions", ACCENT_GREEN, [
            "Apply active promo codes generated by Admin.",
            "Real-time validation against minimum order values.",
            "Instant percentage and flat discount deduction."
        ])
    ]

    for i, (title, sub, col_color, items) in enumerate(user_cards):
        left = Inches(0.8 + i * 3.9)
        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(5.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = CARD_BORDER

        t = slide4.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.3), Inches(4.8))
        tf = t.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = col_color

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_after = Pt(14)

        for item in items:
            pi = tf.add_paragraph()
            pi.text = "✔ " + item
            pi.font.size = Pt(12)
            pi.font.color.rgb = TEXT_MUTED
            pi.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 5: User Module - KYC, Legal Deed & Bookings
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "User Module: Verification, Deed & Booking Lifecycle", "USER EXPERIENCE LAYER")

    cards5 = [
        ("Digital KYC Upload", "Identity Verification", ACCENT_PURPLE, [
            "Upload government-approved Driving License.",
            "Secure image storage via Cloudinary CDN.",
            "Admin-accessible document review pipeline."
        ]),
        ("Rental Agreement Deed", "Legal Terms Compliance", ACCENT_CYAN, [
            "In-app legal terms & conditions agreement.",
            "Covers accident liability, fuel policies, and traffic challans.",
            "Mandatory user acknowledgment before payment."
        ]),
        ("My Bookings Dashboard", "Live Trip Management", ACCENT_GREEN, [
            "Track active, pending, and completed bookings.",
            "View unique Booking ID and digital receipt.",
            "Cancellation support within permitted time window."
        ])
    ]

    for i, (title, sub, col_color, items) in enumerate(cards5):
        left = Inches(0.8 + i * 3.9)
        c = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(5.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = CARD_BORDER

        t = slide5.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.3), Inches(4.8))
        tf = t.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = col_color

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_after = Pt(14)

        for item in items:
            pi = tf.add_paragraph()
            pi.text = "✔ " + item
            pi.font.size = Pt(12)
            pi.font.color.rgb = TEXT_MUTED
            pi.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 6: Admin Module - Command Center & Fleet Control
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "Admin Module: Command Dashboard & Fleet Controls", "ADMIN CONTROL ROOM")

    cards6 = [
        ("Executive Metrics Desk", "Real-Time Operational Pulse", ACCENT_BLUE, [
            "Total vehicle inventory count across categories.",
            "Today's active bookings and confirmed trip counters.",
            "Registered active users and pending booking alerts.",
            "Recent operational audit log stream."
        ]),
        ("Fleet Inventory Manager", "Vehicle Management", ACCENT_CYAN, [
            "Add, update, or remove Bikes, Scooters, and Cars.",
            "Modify hourly rental rates and specifications.",
            "Toggle availability switch (Active / Maintenance).",
            "Direct Cloudinary image uploads."
        ]),
        ("Bookings Management", "Reservation Oversight", ACCENT_GREEN, [
            "Comprehensive tabular list of all user bookings.",
            "Filter bookings by status (Confirmed, Pending, Cancelled).",
            "Inspect customer contact details and trip durations.",
            "One-click booking approval and status management."
        ])
    ]

    for i, (title, sub, col_color, items) in enumerate(cards6):
        left = Inches(0.8 + i * 3.9)
        c = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(5.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = CARD_BORDER

        t = slide6.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.3), Inches(4.8))
        tf = t.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = col_color

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_after = Pt(14)

        for item in items:
            pi = tf.add_paragraph()
            pi.text = "• " + item
            pi.font.size = Pt(12)
            pi.font.color.rgb = TEXT_MUTED
            pi.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 7: Admin Module - User Moderation, Offers & Policies
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "Admin Module: User Moderation, Offers & Governance", "ADMIN CONTROL ROOM")

    cards7 = [
        ("User Management Desk", "Account Governance & Moderation", ACCENT_PURPLE, [
            "Directory of all registered customers with email & phone.",
            "Inspect individual booking history and activity.",
            "Block/Unblock accounts for policy violations.",
            "Password reset and profile verification support."
        ]),
        ("Manage Offers & Coupons", "Promotional Campaign Engine", ACCENT_GREEN, [
            "Create custom promo codes with percentage discounts.",
            "Set validity expiry dates and minimum booking thresholds.",
            "Toggle offer activation status in real time.",
            "Track offer redemption metrics across bookings."
        ]),
        ("Policies & Support Issues", "Customer Care & Compliance", ACCENT_CYAN, [
            "Manage RentHub standard Terms & Conditions Deed.",
            "Support Issues ticketing system with resolution statuses.",
            "Handle customer breakdown and refund requests.",
            "Complete audit trail of support communications."
        ])
    ]

    for i, (title, sub, col_color, items) in enumerate(cards7):
        left = Inches(0.8 + i * 3.9)
        c = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(5.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = CARD_BORDER

        t = slide7.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.3), Inches(4.8))
        tf = t.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = col_color

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_after = Pt(14)

        for item in items:
            pi = tf.add_paragraph()
            pi.text = "• " + item
            pi.font.size = Pt(12)
            pi.font.color.rgb = TEXT_MUTED
            pi.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 8: Enterprise Architecture - Modular In-Development Suite
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "Enterprise Architecture: Planned Growth Modules", "ROADMAP & EXTENSIBILITY")

    grid_items = [
        ("Analytics & Reports", "fas fa-chart-pie", ACCENT_CYAN, "Advanced fleet performance analytics, revenue graphs, and CSV report export engine."),
        ("Sponsor Reports", "fas fa-chart-line", ACCENT_GREEN, "Automated 70/30 revenue distribution ledger and vehicle asset ROI tracking for sponsors."),
        ("Sponsor Portal", "fas fa-handshake", ACCENT_BLUE, "Dedicated partner node interface for live vehicle telemetry and fleet maintenance status."),
        ("Withdrawals Gateway", "fas fa-money-bill-wave", ACCENT_PURPLE, "Multi-bank RTGS/NEFT batch disbursement and UPI settlement gateway.")
    ]

    for idx, (m_title, icon, col, desc) in enumerate(grid_items):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.9)
        top = Inches(1.6 + row * 2.7)

        c = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.4))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col
        c.line.width = Pt(1.2)

        t = slide8.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), Inches(5.1), Inches(2.0))
        tf = t.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = m_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col

        p_status = tf.add_paragraph()
        p_status.text = "Status: Architecture Specified • In Pipeline"
        p_status.font.size = Pt(10)
        p_status.font.bold = True
        p_status.font.color.rgb = TEXT_MUTED
        p_status.space_after = Pt(8)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------
    # SLIDE 9: Security, Authentication & Data Validation
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_header(slide9, "Security, Authentication & Data Integrity", "SECURITY ARCHITECTURE")

    sec_cards = [
        ("JWT Authentication & RBAC", "Access Control", ACCENT_BLUE, [
            "Stateless token authentication for secure sessions.",
            "Strict Role-Based separation between Customers and Admins.",
            "Protected API routes rejecting unauthorized requests."
        ]),
        ("Password & Data Encryption", "Cryptographic Standards", ACCENT_PURPLE, [
            "Passwords hashed using Bcrypt with salt rounds.",
            "Sensitive credentials never stored in plain text.",
            "Sanitized SQL queries preventing injection attacks."
        ]),
        ("Database Relational Integrity", "PostgreSQL Constraints", ACCENT_GREEN, [
            "Foreign keys enforcing referential booking integrity.",
            "Validation preventing concurrent double-booking of vehicles.",
            "Automated timestamping and audit records."
        ])
    ]

    for i, (title, sub, col_color, items) in enumerate(sec_cards):
        left = Inches(0.8 + i * 3.9)
        c = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(5.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = CARD_BORDER

        t = slide9.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.3), Inches(4.8))
        tf = t.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = col_color

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_after = Pt(14)

        for item in items:
            pi = tf.add_paragraph()
            pi.text = "🔒 " + item
            pi.font.size = Pt(12)
            pi.font.color.rgb = TEXT_MUTED
            pi.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 10: Step-by-Step Demonstration Flow
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10)
    add_header(slide10, "Live Demonstration Walkthrough", "PROJECT DEMO PLAN")

    steps = [
        ("Step 1: Customer Booking Flow", ACCENT_CYAN, [
            "Customer registers / logs in to RentHub.",
            "Browses vehicle catalog (Bikes / Scooters / Cars).",
            "Selects pickup duration, applies promo code, and completes booking.",
            "Receives unique Booking ID in 'My Bookings'."
        ]),
        ("Step 2: Admin Command Flow", ACCENT_BLUE, [
            "Admin logs in to the Admin Control Room.",
            "Observes dashboard metrics updating dynamically.",
            "Reviews, verifies, and confirms the customer's booking.",
            "Demonstrates Vehicle Manager, Offers Creator, and Support Desk."
        ])
    ]

    for i, (stitle, scol, sitems) in enumerate(steps):
        left = Inches(0.8 + i * 5.9)
        c = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(5.6), Inches(5.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = scol

        t = slide10.shapes.add_textbox(left + Inches(0.3), Inches(1.9), Inches(5.0), Inches(4.6))
        tf = t.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = scol
        p.space_after = Pt(16)

        for sitem in sitems:
            pi = tf.add_paragraph()
            pi.text = "👉 " + sitem
            pi.font.size = Pt(13)
            pi.font.color.rgb = TEXT_WHITE
            pi.space_after = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 11: Conclusion & Future Scope
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    add_bg(slide11)
    add_header(slide11, "Conclusion & Future Enhancements", "SUMMARY & ROADMAP")

    # Conclusion card
    c_con = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    c_con.fill.solid()
    c_con.fill.fore_color.rgb = CARD_BG
    c_con.line.color.rgb = ACCENT_GREEN

    t_con = slide11.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.6))
    tf_con = t_con.text_frame
    tf_con.word_wrap = True
    p = tf_con.paragraphs[0]
    p.text = "Project Conclusion"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(14)

    con_items = [
        "Delivered a responsive, production-ready vehicle rental web application.",
        "Successfully decoupled User booking experience from Admin operations.",
        "Eliminated physical paperwork through digital agreements & cloud records.",
        "Demonstrated strong software engineering practices with React, Express & Supabase."
    ]
    for ci in con_items:
        p = tf_con.add_paragraph()
        p.text = "✔ " + ci
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # Future scope card
    c_fut = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    c_fut.fill.solid()
    c_fut.fill.fore_color.rgb = CARD_BG
    c_fut.line.color.rgb = ACCENT_PURPLE

    t_fut = slide11.shapes.add_textbox(Inches(7.2), Inches(1.9), Inches(5.0), Inches(4.6))
    tf_fut = t_fut.text_frame
    tf_fut.word_wrap = True
    p = tf_fut.paragraphs[0]
    p.text = "Future Scope"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    p.space_after = Pt(14)

    fut_items = [
        "Native Android & iOS Mobile Applications (React Native).",
        "Hardware IoT GPS tracking & remote keyless engine ignition via OBD-II.",
        "AI-Powered dynamic surge pricing based on weekend & festive demand.",
        "Automated digital damage detection via computer vision snapshot scans."
    ]
    for fi in fut_items:
        p = tf_fut.add_paragraph()
        p.text = "🚀 " + fi
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 12: Thank You / Q&A Slide
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    add_bg(slide12)

    c12 = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(1.2), Inches(9.7), Inches(5.1))
    c12.fill.solid()
    c12.fill.fore_color.rgb = CARD_BG
    c12.line.color.rgb = ACCENT_CYAN
    c12.line.width = Pt(1.5)

    t12 = slide12.shapes.add_textbox(Inches(2.2), Inches(1.6), Inches(8.9), Inches(4.3))
    tf12 = t12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    p_sub = tf12.add_paragraph()
    p_sub.text = "Questions & Evaluator Feedback are Welcome"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = ACCENT_CYAN
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.space_after = Pt(28)

    p_rep = tf12.add_paragraph()
    p_rep.text = "GitHub Repository: https://github.com/Kaushik1575/RentHubR"
    p_rep.font.size = Pt(13)
    p_rep.font.color.rgb = TEXT_MUTED
    p_rep.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "RentHub_Minor_Project_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
